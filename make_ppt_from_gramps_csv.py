from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR 
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.line import LineFormat
from pptx.enum.text import PP_ALIGN
import datetime
import csv
import pandas as pd

# Define dictionnaries for the different csv file sections
Place_Dic={}
Person_Dic={}
Marriage_Dic={}
Family_Dic={}
# Put them in a table
table=[Place_Dic, Person_Dic, Marriage_Dic, Family_Dic]

from cStringIO import StringIO
# read_csv code was copied from https://stackoverflow.com/questions/36904691/pd-read-csv-multiple-tables-and-parse-data-frames-using-index-0/37332651
# thanks for the code!
def read_csv(filename):
	subfiles = [StringIO()]
	with open(filename) as bigfile:
		for line in bigfile:
			if line.strip() == "": # blank line, new subfile
				subfiles.append(StringIO())
			else: # continuation of same subfile                                                                                                                                                   
				if line.strip() != "":
					subfiles[-1].write(line)
			
	i=0		
	for subfile in subfiles:
		subfile.seek(0)
		if subfile.read().strip()!="":
			subfile.seek(0)
			table[i] = pd.read_csv(subfile, sep=',')
			
			#print '*****************'
			#print table[i]
			i=i+1

	bigfile.close()
	# Structure of the GRAMPS CSV file is:
	# # Place feeds with Place,Title,Name,Type,Latitude,Longitude,Code,Enclosed_by,Date
	# # Person,Surname,Given,Call,Suffix,Prefix,Title,Gender,Birth date,Birth place,Birth source,Baptism date,Baptism place,Baptism source,Death date,Death place,Death source,Burial date,Burial place,Burial source,Note
	# # Marriage,Husband,Wife,Date,Place,Source,Note
	# # Family,Child
 
# Get the mother of an Individual parson the marriage he/she is a child of	
# input: '[I-----] person descriptor
# output: the mother person record in the Person_Dic 		
def get_mother(individual):

	return table[1][table[1].Person == table[2][table[2].Marriage == table[3][table[3].Child == individual]['Family'].str.cat()]['Wife'].str.cat()]

# Get the mother of an Individual parson the marriage he/she is a child of
# input: '[I-----] person descriptor
# output: the father person record in the Person_Dic 			
def get_father(individual):

	return table[1][table[1].Person == table[2][table[2].Marriage == table[3][table[3].Child == individual]['Family'].str.cat()]['Husband'].str.cat()]

# Get the record of an Individual parson from his/her person descriptor
# input: '[I-----] person descriptor
# output: the person record in the Person_Dic 	
def get_individual(individual):
	return table[1][table[1].Person == individual]	
	
# Check if an individual is known or not.
# input: '[I-----] person descriptor
# output: True if the person record has a Surname filled and different from Unknown
def get_individual_is_known(individual):
	return (get_individual(individual.Person.str.cat()).Surname.str.cat()!='Unknown' and get_individual(individual.Person.str.cat()).Surname.str.cat()!='' )

# Convert the GRAMPS date format YYYY-MM-DD into DD-MM-YYYY 
def get_date(date):
	if (date !=''):
		return date[8:]+'-'+date[5:7]+'-'+date[:4]
	else:
		return "..."
# Return the birthplace of an individual
# input: '[I-----] person descriptor
# output: return the birth place if filled or  "..." if not.
def get_birth_place(place):
	if table[0][table[0].Place == place].Name.str.cat() != '':
		return table[0][table[0].Place == place].Name.str.cat()
	else:
		return "..."

# Format in a created shape the Surname, First names, Date of birth + ~ + Date of death and place of birth
def print_individual(individual, shape):
	p = shape.text_frame.paragraphs[0]
	# first paragraph
	p.text = individual.Surname.str.cat()+', '+individual.Given.str.cat()

	p.font.name = 'Calibri'
	p.font.size = Pt(14)
	p.font.bold = True	
	p.alignment = PP_ALIGN.CENTER
	
	# second paragraph
	p2 = shape.text_frame.add_paragraph()
	p2.text = get_date(individual['Birth date'].str.cat())+' ~ '+get_date(individual['Death date'].str.cat())+'\n from '+get_birth_place(individual['Birth place'].str.cat())
	p2.font.name = 'Calibri'
	p2.font.size = Pt(10)
	p2.font.bold = True
	p2.alignment = PP_ALIGN.CENTER
	shape.text_frame.word_wrap = False
	shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	shape.text_frame.margin_left=Inches(0)
	shape.text_frame.margin_right=Inches(0)

# add connector between 2 leaves of the family tree
# input: 	shapes as the shape hierarchy of the ppt slide
#			shape_m : parent shape (has to be created before)
#			shape_s : son shape (has to be created before)	
def add_connector(shapes, shape_m, shape_s):
	connector_m = shapes.add_connector(
		MSO_CONNECTOR.CURVE, Inches(2), Inches(2), Inches(1), Inches(1)
	)
	connector_m.begin_connect(shape_m,2)
	connector_m.end_connect(shape_s,0)
	line = connector_m.line
	line.fill.solid()
	line.width=Inches(0.1)
	line.fill.fore_color.rgb = RGBColor(0x96, 0x4B, 0)

# add shape of size and position provided in argument
# input: 	shapes as the shape hierarchy of the ppt slide
#			left_m: x coordinate of left corner of the shape
# 			top_m: y coordinate of top corner of the shape
#			width_m: width of the shape
#			height_m: height of the shape
# output:	the created parent shape
def add_parent_shape(shapes, left_m, top_m, width_m, height_m):
	shape_m = shapes.add_shape(
		MSO_SHAPE.ROUNDED_RECTANGLE, left_m, top_m, width_m, height_m
	)
	fill = shape_m.fill
	fill.solid()
	fore_color = fill.fore_color
	fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
	shape_m.line.fill.solid()
	shape_m.line.fill.fore_color.rgb = RGBColor(107,142,35)
	return shape_m

# add parent shapes of an individual for the current level in the tree among total level (needed to space horizontally) 
# input: 	shapes as the shape hierarchy of the ppt slide
#			shape_s: shape of the child (must be created before)
#			current_level: current level in the tree. From 1 to total_level
#			total_level: total number of levels needed in the tree.
# 			'[I-----]' person descriptor to add parents to
# output:	return the mother and father shapes
def add_parents(shapes, shape_s, current_level, total_level, individual):

	top_m = shape_s.top-2*shape_s.height

	# mother
	left_m = shape_s.left - 2**(total_level-current_level)*shape_s.width/2

	# father
	left_f = shape_s.left + 2**(total_level-current_level)*shape_s.width/2

	shape_m=add_parent_shape(shapes, left_m, top_m, shape_s.width, shape_s.height)
	shape_f=add_parent_shape(shapes, left_f, top_m, shape_s.width, shape_s.height)
	
	add_connector(shapes, shape_m, shape_s)
	add_connector(shapes, shape_f, shape_s)
	
	husband = get_father(individual)
	wife = get_mother(individual)

	print_individual(husband, shape_m)
	print_individual(wife, shape_f)

	
	return [shape_m,shape_f]
	
# process recursively the family tree sill no ancestor is found 
# input: 	shapes as the shape hierarchy of the ppt slide
#			father: shape of the father
#			mother: shape of the mother
#			current: current level in the tree
#			total: total level of the tree
#			husband: husband '[I-----]' person descriptor
#			wife: wife '[I-----]' person descriptor
def process (shapes, father, mother, current, total, husband, wife):
	if current == total:
		return
	else:
		if (get_individual_is_known(husband)):
			[father2, mother2] = add_parents(shapes, father, current, total, husband.Person.str.cat())
			process (shapes, father2, mother2, current+1, total, get_father(husband.Person.str.cat()), get_mother(husband.Person.str.cat()))

		if (get_individual_is_known(wife)):
			[father3, mother3] = add_parents(shapes, mother, current, total, wife.Person.str.cat())
			process (shapes, father3, mother3, current+1, total, get_father(wife.Person.str.cat()), get_mother(wife.Person.str.cat()))

			
#read csv file			
read_csv('C:\Users\Philippe\Documents\Untitled_1.csv')
		
# Initialization 
individual = get_individual('[I0001]')
max = 6
	
# Create the blank ppt	
prs = Presentation()
title_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_slide_layout)

# Get the blank ppt dimensions
s_height = prs.slide_height
s_width = prs.slide_width


# Create the initial child shape
shapes = slide.shapes
width_s = Inches(1.0)
height_s = Inches(0.5)
left_s = s_width/2-width_s/2
top_s = s_height - height_s

shape_s = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left_s, top_s, width_s, height_s
)
# Populate name and dates
p = shape_s.text_frame.paragraphs[0]
run = p.add_run()
run.text = individual.Surname.str.cat()+', '+individual.Given.str.cat()+'\n'+get_date(individual['Birth date'].str.cat())+' ~ '+get_date(individual['Death date'].str.cat())+'\n de '+get_birth_place(individual['Birth place'].str.cat())

font = run.font
font.name = 'Calibri'
font.size = Pt(12)
font.bold = True

shape_s.text_frame.vertical_anchor = MSO_ANCHOR.TOP
shape_s.text_frame.word_wrap = False
shape_s.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
shape_s.text_frame.margin_left=Inches(0)
shape_s.text_frame.margin_right=Inches(0)
fill = shape_s.fill
fill.solid()
fore_color = fill.fore_color

fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
shape_s.line.fill.solid()
shape_s.line.fill.fore_color.rgb = RGBColor(107,142,35)



# create the first node
[father, mother] = add_parents(shapes, shape_s, 1,max, individual.Person.str.cat())
# get the husband and wife of the first marriage above
husband = get_father(individual.Person.str.cat())
wife = get_mother(individual.Person.str.cat())
# enter the recursive family tree processing
process (shapes, father, mother, 2, max, husband, wife )

# save the generated pptx
prs.save('test.pptx')